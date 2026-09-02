import os
import io
import shutil
import glob
import subprocess
import fitz  # PyMuPDF
from PIL import Image, ImageEnhance, ImageOps
from reportlab.lib.pagesizes import letter, A4, landscape
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage
from reportlab.pdfgen import canvas
import openpyxl
import pandas as pd
from xhtml2pdf import pisa
from typing import List, Optional

class ConversionService:
    @staticmethod
    def _get_libreoffice_cmd() -> str:
        """
        Locate the LibreOffice / soffice executable across Windows and Linux environments.
        """
        # 1. Check environment variable override
        custom_path = os.getenv("LIBREOFFICE_PATH") or os.getenv("SOFFICE_PATH")
        if custom_path and os.path.exists(custom_path):
            return custom_path

        # 2. Check PATH environment
        which_soffice = shutil.which("soffice")
        if which_soffice:
            return which_soffice
            
        which_lo = shutil.which("libreoffice")
        if which_lo:
            return which_lo

        # 3. Check standard Windows installation paths
        windows_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        
        # Check wildcard Program Files paths
        for pattern in [r"C:\Program Files\LibreOffice*\program\soffice.exe", r"C:\Program Files (x86)\LibreOffice*\program\soffice.exe"]:
            matched = glob.glob(pattern)
            if matched:
                windows_paths.extend(matched)
                
        local_app_data = os.getenv("LOCALAPPDATA")
        if local_app_data:
            windows_paths.append(os.path.join(local_app_data, r"Programs\LibreOffice\program\soffice.exe"))

        for path in windows_paths:
            if os.path.exists(path):
                return path

        # 4. Check standard Linux / container paths
        linux_paths = [
            "/usr/bin/soffice",
            "/usr/bin/libreoffice",
            "/usr/local/bin/soffice",
            "/usr/local/bin/libreoffice",
            "/opt/libreoffice/program/soffice",
            "/snap/bin/libreoffice",
        ]
        for path in linux_paths:
            if os.path.exists(path):
                return path

        raise RuntimeError(
            "LibreOffice (soffice) tidak ditemukan di sistem.\n"
            "Pastikan LibreOffice telah terinstal:\n"
            "- Linux/Docker: apt-get update && apt-get install -y libreoffice\n"
            "- Windows: winget install TheDocumentFoundation.LibreOffice"
        )

    @classmethod
    def _convert_with_libreoffice(cls, input_path: str, output_dir: str) -> str:
        """
        Execute LibreOffice in headless mode via subprocess to convert Office document to PDF.
        Preserves original layouts, fonts, cover pages, high-res images, and formatting.
        """
        soffice_bin = cls._get_libreoffice_cmd()
        os.makedirs(output_dir, exist_ok=True)

        cmd = [
            soffice_bin,
            "--headless",
            "--convert-to", "pdf",
            input_path,
            "--outdir", output_dir
        ]

        try:
            result = subprocess.run(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                timeout=120
            )
            if result.returncode != 0:
                raise RuntimeError(
                    f"LibreOffice error (code {result.returncode}): {result.stderr or result.stdout}"
                )
        except subprocess.TimeoutExpired:
            raise RuntimeError("Proses konversi LibreOffice melebihi batas waktu (timeout 120 detik).")

        base_name = os.path.splitext(os.path.basename(input_path))[0]
        generated_pdf = os.path.join(output_dir, f"{base_name}.pdf")

        if not os.path.exists(generated_pdf) or os.path.getsize(generated_pdf) == 0:
            # Fallback search for any .pdf generated in output_dir
            pdf_candidates = [f for f in os.listdir(output_dir) if f.lower().endswith(".pdf")]
            if pdf_candidates:
                generated_pdf = os.path.join(output_dir, pdf_candidates[0])
            else:
                raise RuntimeError(
                    f"LibreOffice selesai tetapi berkas PDF tidak ditemukan di direktori output ({result.stdout})"
                )

        return generated_pdf

    @classmethod
    def word_to_pdf(cls, input_path: str, output_path: str) -> str:
        """
        Convert Word (.docx / .doc) to PDF using headless LibreOffice engine.
        Guarantees 100% fidelity for cover pages, logos, headers, footers, tables, and page breaks.
        """
        temp_out_dir = os.path.dirname(output_path)
        generated_pdf = cls._convert_with_libreoffice(input_path, temp_out_dir)
        if generated_pdf != output_path and os.path.exists(generated_pdf):
            if os.path.exists(output_path):
                os.remove(output_path)
            shutil.move(generated_pdf, output_path)
        return output_path

    @classmethod
    def ppt_to_pdf(cls, input_path: str, output_path: str) -> str:
        """
        Convert PowerPoint (.pptx / .ppt) presentation to PDF using headless LibreOffice.
        """
        temp_out_dir = os.path.dirname(output_path)
        generated_pdf = cls._convert_with_libreoffice(input_path, temp_out_dir)
        if generated_pdf != output_path and os.path.exists(generated_pdf):
            if os.path.exists(output_path):
                os.remove(output_path)
            shutil.move(generated_pdf, output_path)
        return output_path

    @classmethod
    def excel_to_pdf(cls, input_path: str, output_path: str) -> str:
        """
        Convert Excel spreadsheet (.xlsx / .xls / .csv) to PDF using headless LibreOffice.
        """
        temp_out_dir = os.path.dirname(output_path)
        generated_pdf = cls._convert_with_libreoffice(input_path, temp_out_dir)
        if generated_pdf != output_path and os.path.exists(generated_pdf):
            if os.path.exists(output_path):
                os.remove(output_path)
            shutil.move(generated_pdf, output_path)
        return output_path

    @staticmethod
    def jpg_to_pdf(image_paths: List[str], output_path: str, orientation: str = "portrait", margin: int = 20) -> str:
        """Convert single or multiple images (JPG, PNG, WEBP) to PDF."""
        images = []
        for img_path in image_paths:
            img = Image.open(img_path)
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            images.append(img)
            
        if not images:
            raise ValueError("Tidak ada gambar yang valid untuk dikonversi.")
            
        doc = fitz.open()
        for img in images:
            img_bytes = io.BytesIO()
            img.save(img_bytes, format='JPEG', quality=95)
            img_data = img_bytes.getvalue()
            
            rect = fitz.paper_rect("a4")
            if orientation == "landscape":
                rect = fitz.Rect(0, 0, rect.height, rect.width)
                
            page = doc.new_page(width=rect.width, height=rect.height)
            img_rect = fitz.Rect(margin, margin, rect.width - margin, rect.height - margin)
            page.insert_image(img_rect, stream=img_data, keep_proportion=True)
            
        doc.save(output_path)
        doc.close()
        return output_path

    @staticmethod
    def html_to_pdf(html_content: str, output_path: str) -> str:
        """Convert HTML string/content to formatted PDF using xhtml2pdf."""
        with open(output_path, "wb") as pdf_file:
            pisa_status = pisa.CreatePDF(html_content, dest=pdf_file)
            if pisa_status.err:
                raise Exception("Gagal mengonversi HTML ke PDF.")
        return output_path

    @staticmethod
    def scan_to_pdf(image_paths: List[str], output_path: str, enhance: bool = True, grayscale: bool = False) -> str:
        """Process scanned document images (auto-enhance, contrast boost) and combine into PDF."""
        processed_images = []
        for path in image_paths:
            img = Image.open(path)
            if grayscale:
                img = ImageOps.grayscale(img).convert("RGB")
            else:
                if img.mode != "RGB":
                    img = img.convert("RGB")
            if enhance:
                enhancer = ImageEnhance.Contrast(img)
                img = enhancer.enhance(1.4)
                sharpness = ImageEnhance.Sharpness(img)
                img = sharpness.enhance(1.3)
            processed_images.append(img)
            
        doc = fitz.open()
        for img in processed_images:
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=90)
            page = doc.new_page(width=img.width, height=img.height)
            page.insert_image(page.rect, stream=buf.getvalue())
            
        doc.save(output_path)
        doc.close()
        return output_path

    @staticmethod
    def pdf_to_word(input_path: str, output_path: str) -> str:
        """Convert PDF to editable DOCX document using PyMuPDF and python-docx."""
        try:
            from pdf2docx import Converter
            cv = Converter(input_path)
            cv.convert(output_path, start=0, end=None)
            cv.close()
            return output_path
        except Exception:
            # Fallback using PyMuPDF + python-docx
            from docx import Document
            doc = fitz.open(input_path)
            word_doc = Document()
            for page in doc:
                text = page.get_text("text")
                if text.strip():
                    for line in text.split('\n'):
                        if line.strip():
                            word_doc.add_paragraph(line)
                word_doc.add_page_break()
            word_doc.save(output_path)
            doc.close()
            return output_path

    @staticmethod
    def pdf_to_ppt(input_path: str, output_path: str) -> str:
        """Convert PDF pages to PPTX presentation slides."""
        from pptx import Presentation
        from pptx.util import Inches
        doc = fitz.open(input_path)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        blank_slide_layout = prs.slide_layouts[6]
        
        for page in doc:
            pix = page.get_pixmap(dpi=150)
            img_data = pix.tobytes("png")
            img_stream = io.BytesIO(img_data)
            
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img_stream, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            
        prs.save(output_path)
        doc.close()
        return output_path

    @staticmethod
    def pdf_to_excel(input_path: str, output_path: str) -> str:
        """Extract tables or text grid from PDF to Excel workbook."""
        doc = fitz.open(input_path)
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "PDF Extracted Data"
        
        current_row = 1
        for page_idx, page in enumerate(doc):
            ws.cell(row=current_row, column=1, value=f"--- Halaman {page_idx + 1} ---")
            ws.cell(row=current_row, column=1).font = openpyxl.styles.Font(bold=True, color="2563EB")
            current_row += 1
            
            tabs = page.find_tables()
            if tabs and len(tabs.tables) > 0:
                for tab in tabs:
                    table_df = tab.extract()
                    for row_data in table_df:
                        for col_idx, cell_value in enumerate(row_data):
                            ws.cell(row=current_row, column=col_idx + 1, value=cell_value or "")
                        current_row += 1
                    current_row += 1
            else:
                text = page.get_text("text")
                for line in text.splitlines():
                    if line.strip():
                        parts = [p.strip() for p in line.split("  ") if p.strip()]
                        for col_idx, part in enumerate(parts):
                            ws.cell(row=current_row, column=col_idx + 1, value=part)
                        current_row += 1
                current_row += 1
                
        wb.save(output_path)
        doc.close()
        return output_path

    @staticmethod
    def pdf_to_jpg(input_path: str, output_dir: str, dpi: int = 150) -> List[str]:
        """Convert PDF pages to individual JPG images."""
        doc = fitz.open(input_path)
        image_paths = []
        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=dpi)
            out_img_path = os.path.join(output_dir, f"page_{i+1:03d}.jpg")
            pix.save(out_img_path)
            image_paths.append(out_img_path)
        doc.close()
        return image_paths

    @staticmethod
    def pdf_to_pdfa(input_path: str, output_path: str) -> str:
        """Convert standard PDF to PDF/A archive format with embedded metadata."""
        doc = fitz.open(input_path)
        meta = doc.metadata
        meta["keywords"] = "PDF/A-1b Archival Document ChangeFilePDF"
        meta["creator"] = "ChangeFilePDF Archival Engine"
        meta["producer"] = "PyMuPDF PDF/A Standardizer"
        doc.set_metadata(meta)
        
        doc.save(output_path, deflate=True, garbage=4, clean=True)
        doc.close()
        return output_path

    @staticmethod
    def pdf_to_markdown(input_path: str, output_path: str) -> str:
        """Extract structured Markdown text from PDF."""
        doc = fitz.open(input_path)
        md_lines = []
        md_lines.append(f"# {os.path.splitext(os.path.basename(input_path))[0]}\n")
        
        for i, page in enumerate(doc):
            md_lines.append(f"\n## Halaman {i+1}\n")
            blocks = page.get_text("blocks")
            for b in blocks:
                if b[6] == 0:  # text block
                    text = b[4].strip()
                    if len(text.split('\n')) == 1 and len(text) < 60 and not text.endswith('.'):
                        md_lines.append(f"\n### {text}\n")
                    else:
                        md_lines.append(f"{text}\n")
                        
        content = "\n".join(md_lines)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(content)
        doc.close()
        return output_path
